"""
Template analyzer for target documents.
Uses section-based analysis: identifies Heading + Body pairs to create semantic sections.
"""
import logging
import re
from pathlib import Path
from typing import Literal, Any

from pydantic import BaseModel

from app.core.exceptions import (
    ParsingError, 
    UnsupportedFileTypeError
)

logger = logging.getLogger(__name__)


class TemplateSection(BaseModel):
    """A section in the template: heading + body content."""
    section_id: str
    heading_text: str
    heading_paragraph_idx: int  # Index of heading paragraph
    body_start_idx: int  # First body paragraph index
    body_end_idx: int  # Last body paragraph index (exclusive)
    body_preview: str  # First ~100 chars of body for context
    section_type: Literal["title", "section", "subsection"]


class TemplateDNA(BaseModel):
    """Template styling DNA for reconstruction."""
    heading_style_name: str
    heading_font_name: str | None = None
    heading_font_size: int | None = None  # In points
    heading_font_color: str | None = None  # Hex color
    heading_font_bold: bool | None = None
    
    # Subheading (Heading 2) style
    subheading_style_name: str = "Heading 2"
    subheading_font_name: str | None = None
    subheading_font_size: int | None = None
    subheading_font_color: str | None = None
    subheading_font_bold: bool | None = None
    
    body_style_name: str
    body_font_name: str | None = None
    body_font_size: int | None = None  # In points
    body_font_color: str | None = None
    body_font_bold: bool | None = None
    body_font_italic: bool | None = None
    
    # Bullet point style
    bullet_style_name: str = "List Bullet"
    bullet_font_name: str | None = None
    bullet_font_size: int | None = None
    
    safe_zone_end_idx: int  # Everything before this index is preserved (cover, TOC)
    first_content_section_idx: int  # Index of first replaceable section


class TemplateAnalysis(BaseModel):
    """Result of template section analysis."""
    sections: list[TemplateSection]
    section_ids: list[str]
    template_file: str
    total_paragraphs: int
    content_paragraph_indices: list[int] = []  # Indices of CONTENT paragraphs to replace
    template_dna: TemplateDNA | None = None  # Template DNA for reconstruction


def analyze_template(file_path: Path) -> TemplateAnalysis:
    """
    Analyze a template document to detect sections (Heading + Body pairs).
    
    Args:
        file_path: Path to the template document
        
    Returns:
        TemplateAnalysis with detected sections
    """
    suffix = file_path.suffix.lower()
    
    if suffix == ".docx":
        return _analyze_docx_sections(file_path)
    elif suffix == ".pptx":
        return _analyze_pptx_sections(file_path)
    elif suffix == ".pdf":
        return _analyze_pdf_sections(file_path)
    else:
        raise UnsupportedFileTypeError(f"Unsupported template type: {suffix}")


def _analyze_docx_sections(file_path: Path) -> TemplateAnalysis:
    """
    Analyze DOCX template - extract template DNA for reconstruction.
    Finds master section styles and safe zones.
    """
    try:
        from docx import Document
        
        doc = Document(str(file_path))
        total_paras = len(doc.paragraphs)
        
        logger.info(f"Analyzing DOCX template: {file_path.name} ({total_paras} paragraphs)")
        
        # Find the first major heading (Heading 1) - this is our master section
        first_heading_idx = -1
        heading_style_name = "Heading 1"
        heading_font_name = None
        heading_font_size = None
        heading_font_color = None
        heading_font_bold = None
        
        body_style_name = "Normal"
        body_font_name = None
        body_font_size = None
        body_font_color = None
        body_font_bold = None
        body_font_italic = None
        safe_zone_end = 0
        
        for idx, para in enumerate(doc.paragraphs):
            text = para.text.strip()
            if not text:
                continue
            
            style_name = para.style.name if para.style else "Normal"
            style_lower = style_name.lower()
            
            # Look for first Heading 1
            if "heading 1" in style_lower and first_heading_idx == -1:
                first_heading_idx = idx
                heading_style_name = style_name
                safe_zone_end = idx  # Everything before this is the safe zone
                
                # Extract font properties from heading
                if para.runs:
                    first_run = para.runs[0]
                    heading_font_name = first_run.font.name
                    if first_run.font.size:
                        heading_font_size = int(first_run.font.size.pt)
                    heading_font_bold = first_run.font.bold
                    # Try to get color
                    if first_run.font.color and first_run.font.color.rgb:
                        heading_font_color = str(first_run.font.color.rgb)
                
                logger.info(f"  Found master heading at para {idx}:")
                logger.info(f"    Style: '{style_name}'")
                logger.info(f"    Font: {heading_font_name}, Size: {heading_font_size}, Bold: {heading_font_bold}")
                logger.info(f"    Text: '{text[:40]}'")
                
                # Next non-empty, non-heading paragraph is the body style
                for body_idx in range(idx + 1, min(idx + 10, total_paras)):
                    body_para = doc.paragraphs[body_idx]
                    body_text = body_para.text.strip()
                    if body_text:
                        body_style = body_para.style.name if body_para.style else "Normal"
                        body_style_lower = body_style.lower()
                        if "heading" not in body_style_lower and "toc" not in body_style_lower:
                            body_style_name = body_style
                            
                            # Extract font properties from body
                            if body_para.runs:
                                first_body_run = body_para.runs[0]
                                body_font_name = first_body_run.font.name
                                if first_body_run.font.size:
                                    body_font_size = int(first_body_run.font.size.pt)
                                body_font_bold = first_body_run.font.bold
                                body_font_italic = first_body_run.font.italic
                                if first_body_run.font.color and first_body_run.font.color.rgb:
                                    body_font_color = str(first_body_run.font.color.rgb)
                            
                            logger.info(f"  Found master body at para {body_idx}:")
                            logger.info(f"    Style: '{body_style}'")
                            logger.info(f"    Font: {body_font_name}, Size: {body_font_size}")
                            break
                break
        
        # Extract template DNA
        template_dna = TemplateDNA(
            heading_style_name=heading_style_name,
            heading_font_name=heading_font_name,
            heading_font_size=heading_font_size,
            heading_font_color=heading_font_color,
            heading_font_bold=heading_font_bold,
            body_style_name=body_style_name,
            body_font_name=body_font_name,
            body_font_size=body_font_size,
            body_font_color=body_font_color,
            body_font_bold=body_font_bold,
            body_font_italic=body_font_italic,
            safe_zone_end_idx=safe_zone_end,
            first_content_section_idx=first_heading_idx if first_heading_idx != -1 else 0
        )
        
        logger.info(f"  Template DNA extracted:")
        logger.info(f"    Heading style: {heading_style_name}")
        logger.info(f"    Body style: {body_style_name}")
        logger.info(f"    Safe zone: 0-{safe_zone_end}")
        logger.info(f"    First content section: {first_heading_idx}")
        
        # Create a single section for compatibility
        section = TemplateSection(
            section_id="sec_all",
            heading_text="Document",
            heading_paragraph_idx=-1,
            body_start_idx=0,
            body_end_idx=total_paras,
            body_preview="Template DNA extraction",
            section_type="section"
        )
        
        return TemplateAnalysis(
            sections=[section],
            section_ids=["sec_all"],
            template_file=file_path.name,
            total_paragraphs=total_paras,
            template_dna=template_dna
        )
    except Exception as e:
        raise AnalysisError(f"Failed to analyze DOCX sections: {str(e)}")


def _is_structural_paragraph(para, style_name: str, text: str) -> bool:
    """
    Determine if a paragraph is STRUCTURAL (keep) or CONTENT (replace).
    
    STRUCTURAL: Only paragraphs with images or in header/footer sections
    CONTENT: All text paragraphs (titles, headings, TOC, body text)
    """
    # Check for images or shapes - these are STRUCTURAL
    try:
        # Check if paragraph contains images
        if para._element.xpath('.//pic:pic'):
            return True  # Has images - preserve
    except:
        pass
    
    # Check if in header/footer section (not just "Header" style)
    try:
        # This is more complex - for now, we'll classify by style
        style_lower = style_name.lower()
        if 'header' in style_lower or 'footer' in style_lower:
            # Only if it's an actual header/footer style (not "Heading")
            if 'heading' not in style_lower:
                return True
    except:
        pass
    
    # Everything else is CONTENT (will be replaced)
    return False


def _analyze_pptx_sections(file_path: Path) -> TemplateAnalysis:
    """
    Analyze PPTX template: each slide is a section.
    """
    try:
        from pptx import Presentation
        
        prs = Presentation(str(file_path))
        sections: list[TemplateSection] = []
        section_ids: list[str] = []
        
        for slide_idx, slide in enumerate(prs.slides):
            title_text = f"Slide {slide_idx + 1}"
            body_preview = ""
            
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                
                # Check if title
                is_title = False
                if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                    ph_type = shape.placeholder_format.type
                    if ph_type in [1, 2, 3]:
                        is_title = True
                        title_text = text
                
                if not is_title and not body_preview:
                    body_preview = text[:100]
            
            section_id = f"slide_{slide_idx}"
            sections.append(TemplateSection(
                section_id=section_id,
                heading_text=title_text,
                heading_paragraph_idx=slide_idx,
                body_start_idx=slide_idx,
                body_end_idx=slide_idx + 1,
                body_preview=body_preview,
                section_type="section"
            ))
            section_ids.append(section_id)
        
        return TemplateAnalysis(
            sections=sections,
            section_ids=section_ids,
            template_file=file_path.name,
            total_paragraphs=len(prs.slides)
        )
        
    except Exception as e:
        raise ParsingError(f"Failed to analyze PPTX sections: {str(e)}", details=str(e))


def _analyze_pdf_sections(file_path: Path) -> TemplateAnalysis:
    """
    Analyze PDF template: basic page-based sections.
    Note: PDF modification is limited, this is for read purposes.
    """
    try:
        import fitz
        
        doc = fitz.open(str(file_path))
        sections: list[TemplateSection] = []
        section_ids: list[str] = []
        
        for page_idx, page in enumerate(doc):
            text = page.get_text()
            lines = text.strip().split('\n')
            
            title = lines[0][:50] if lines else f"Page {page_idx + 1}"
            body_preview = ' '.join(lines[1:5])[:100] if len(lines) > 1 else ""
            
            section_id = f"page_{page_idx}"
            sections.append(TemplateSection(
                section_id=section_id,
                heading_text=title,
                heading_paragraph_idx=page_idx,
                body_start_idx=page_idx,
                body_end_idx=page_idx + 1,
                body_preview=body_preview,
                section_type="section"
            ))
            section_ids.append(section_id)
        
        doc.close()
        
        return TemplateAnalysis(
            sections=sections,
            section_ids=section_ids,
            template_file=file_path.name,
            total_paragraphs=len(sections)
        )
        
    except Exception as e:
        raise ParsingError(f"Failed to analyze PDF sections: {str(e)}", details=str(e))


def get_section_descriptions(analysis: TemplateAnalysis) -> str:
    """
    Generate descriptions of template sections for AI prompt.
    """
    lines = []
    for sec in analysis.sections:
        context = f" (Context: '{sec.body_preview[:80]}...')" if sec.body_preview else ""
        lines.append(f"- {sec.section_id}: \"{sec.heading_text}\"{context}")
    
    return "\n".join(lines)
