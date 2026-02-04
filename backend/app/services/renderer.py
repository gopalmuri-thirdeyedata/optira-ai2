"""
Content replacement engine for template documents.
Performs style-preserving section injection: replaces section body content while
preserving the template's heading styles, fonts, header/footer, images, and layout.
"""
import logging
import re
from pathlib import Path
from typing import Any, Optional

from app.core.exceptions import RenderingError, UnsupportedFileTypeError
from app.services.ai_mapper import SectionMapping
from app.services.analyzer import TemplateAnalysis

logger = logging.getLogger(__name__)


def clean_bullet_text(text: str) -> str:
    """
    Remove leading bullet markers from text to prevent duplicates.
    Handles: -, *, •, numbered lists (1., 2.), lettered lists (a., b.)
    """
    # Strip leading whitespace first
    text = text.strip()
    # Remove common bullet/list markers at the start
    patterns = [
        r'^[-*•]\s*',           # Dash, asterisk, bullet
        r'^\d+\.\s*',           # Numbered (1., 2., etc.)
        r'^[a-zA-Z]\.\s*',      # Lettered (a., b., etc.)
        r'^[a-zA-Z]\)\s*',      # Lettered with paren (a), b))
    ]
    for pattern in patterns:
        text = re.sub(pattern, '', text)
    return text.strip()


def _detect_cover_title(doc, safe_zone_end):
    """
    Find the main title paragraph on the cover page using multiple heuristics.
    
    Returns:
        tuple: (paragraph_index, paragraph_object) or (None, None) if not found
    """
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    
    candidates = []
    
    # Search first 10 paragraphs or safe zone, whichever is smaller
    search_range = min(10, safe_zone_end)
    
    for idx in range(search_range):
        para = doc.paragraphs[idx]
        text = para.text.strip()
        
        # Skip empty or very short text
        if not text or len(text) < 5:
            continue
            
        score = 0
        
        # 1. Style-based detection (most reliable)
        style_name = para.style.name.lower() if para.style else ""
        if 'title' in style_name:
            score += 100
        if 'cover' in style_name:
            score += 80
        if 'heading' in style_name and '1' in style_name:
            score += 50
            
        # 2. Formatting-based detection
        if para.runs:
            run = para.runs[0]
            
            # Font size (larger = more likely title)
            if run.font.size:
                font_pt = run.font.size.pt if hasattr(run.font.size, 'pt') else 0
                if font_pt > 0:
                    score += min(font_pt, 50)  # Cap at 50 points
            
            # Bold text
            if run.font.bold:
                score += 20
                
        # 3. Alignment (centered titles are common)
        try:
            if para.alignment == WD_ALIGN_PARAGRAPH.CENTER:
                score += 15
        except:
            pass
            
        # 4. Position bonus (earlier = more likely)
        score += (10 - idx) * 2
        
        # Skip very long text (likely not a title)
        if len(text) > 150:
            score -= 30
            
        candidates.append((score, idx, para, text))
    
    # Return highest scoring candidate
    if candidates:
        candidates.sort(reverse=True, key=lambda x: x[0])
        score, idx, para, text = candidates[0]
        
        # Only return if score is reasonable
        if score > 20:
            logger.info(f"  Detected cover title at para {idx} (score={score}): '{text[:50]}'")
            return idx, para
    
    logger.info("  No clear cover title detected")
    return None, None


def _update_safe_zone(doc, template_dna, section_titles: list[str], document_title: str = None):
    """
    Update the safe zone (cover page, TOC) with new content information.
    
    Args:
        doc: Document object
        template_dna: Template metadata
        section_titles: List of section/chapter titles from AI
        document_title: Optional document title for cover page (defaults to first section if None)
    """
    safe_zone_end = template_dna.safe_zone_end_idx
    
    logger.info(f"  Updating safe zone (cover + TOC)...")
    
    # === STEP 1: Update Cover Page Title ===
    if not document_title and section_titles:
        # Use first section title as document title if not provided
        document_title = section_titles[0]
    
    if document_title:
        title_idx, title_para = _detect_cover_title(doc, safe_zone_end)
        if title_para:
            try:
                # Preserve all formatting, just change the text content
                old_title = title_para.text
                if title_para.runs:
                    # Update first run with new title
                    title_para.runs[0].text = document_title
                    # Clear subsequent runs to avoid leftover text
                    for run in title_para.runs[1:]:
                        run.text = ""
                    logger.info(f"  ✓ Updated cover title: '{old_title[:30]}...' → '{document_title[:30]}'")
                else:
                    # Fallback: add run if none exist
                    title_para.add_run(document_title)
                    logger.info(f"  ✓ Set cover title: '{document_title}'")
            except Exception as e:
                logger.warning(f"  Failed to update cover title: {e}")
    
    # === STEP 2: Find TOC Entries ===
    toc_start = -1
    toc_entries = []
    
    for idx in range(safe_zone_end):
        para = doc.paragraphs[idx]
        text = para.text.strip()
        style_name = para.style.name.lower() if para.style else ""
        
        # Detect TOC heading
        if "table of contents" in text.lower() or "contents" in text.lower():
            toc_start = idx
            logger.info(f"  Found TOC heading at para {idx}: '{text}'")
            continue
        
        # Look for TOC-specific styles
        if "toc" in style_name and idx > toc_start:
            toc_entries.append(idx)
            continue
        
        # If we've found TOC heading, look for entry patterns
        if toc_start >= 0 and text:
            # Pattern 1: Has tab character (common in TOCs)
            has_tab = "\t" in para.text
            
            # Pattern 2: Has dotted leader
            has_dots = "..." in text or "․․․" in text
            
            # Pattern 3: Starts with number (1., 1, etc.)
            has_number = re.match(r'^\d+[\.\)]\s+\w+', text)
            
            # Pattern 4: Ends with number (page number)
            ends_with_number = re.search(r'\d+\s*$', text)
            
            if has_tab or has_dots or has_number or ends_with_number:
                toc_entries.append(idx)
    
    logger.info(f"  Found {len(toc_entries)} TOC entries")
    
    # === STEP 3: Update TOC Entries ===
    if toc_entries and section_titles:
        updated_count = 0
        for i, para_idx in enumerate(toc_entries):
            if i >= len(section_titles):
                # More TOC entries than sections - break
                break
                
            para = doc.paragraphs[para_idx]
            new_title = section_titles[i]
            old_text = para.text
            
            try:
                # Strategy: Replace title portion, keep formatting & trailing elements
                # TOC format examples:
                #   "1. Section Title ........ 5"
                #   "Section Title\t5"
                #   "1. Section Title"
                
                # Find where the title text ends (before dots, tabs, or trailing numbers)
                text = para.text
                
                # Look for separators
                dot_match = re.search(r'\.{2,}', text)  # Multiple dots
                tab_idx = text.find('\t')
                
                # Determine where to split
                split_idx = len(text)
                if dot_match:
                    split_idx = min(split_idx, dot_match.start())
                if tab_idx > 0:
                    split_idx = min(split_idx, tab_idx)
                
                # Preserve trailing part (dots, tabs, page numbers)
                trailing_part = text[split_idx:] if split_idx < len(text) else ""
                
                # Build new TOC entry
                # Remove old numbering if present
                new_title_clean = re.sub(r'^\d+[\.\)]\s*', '', new_title)
                new_text = f"{i+1}. {new_title_clean}{trailing_part}"
                
                # Update paragraph text while preserving formatting
                if para.runs:
                    # Update first run
                    para.runs[0].text = new_text
                    # Clear other runs (except we want to preserve tab/page number runs)
                    # Actually, safer to just update first run with full text
                    for run in para.runs[1:]:
                        run.text = ""
                else:
                    para.add_run(new_text)
                
                logger.info(f"  ✓ TOC entry {i+1}: '{old_text[:40]}'  →  '{new_text[:40]}'")
                updated_count += 1
                
            except Exception as e:
                logger.warning(f"  Failed to update TOC entry {i+1}: {e}")
        
        logger.info(f"  Updated {updated_count}/{len(section_titles)} TOC entries")
    else:
        if not toc_entries:
            logger.info(f"  ⚠ No TOC entries found to update")
        elif not section_titles:
            logger.info(f"  ⚠ No section titles to populate TOC")



def render_document(
    template_path: Path,
    output_path: Path,
    mapping: SectionMapping,
    analysis: TemplateAnalysis
) -> Path:
    """
    Replace section body content in template with mapped content.
    Preserves all template formatting, headers, footers, images.
    
    Args:
        template_path: Path to the template document
        output_path: Path where the rendered document will be saved
        mapping: Section-to-content mapping from AI
        analysis: Template analysis with section info
        
    Returns:
        Path to the rendered document
    """
    suffix = template_path.suffix.lower()
    
    if suffix == ".docx":
        return _render_docx_sections(template_path, output_path, mapping, analysis)
    elif suffix == ".pptx":
        return _render_pptx_sections(template_path, output_path, mapping, analysis)
    elif suffix == ".pdf":
        raise UnsupportedFileTypeError(
            "PDF template replacement is not directly supported. "
            "Consider using a DOCX template and converting to PDF."
        )
    else:
        raise UnsupportedFileTypeError(f"Unsupported template type: {suffix}")


def _render_docx_sections(
    template_path: Path,
    output_path: Path,
    mapping: SectionMapping,
    analysis: TemplateAnalysis
) -> Path:
    """
    Render DOCX using template-driven reconstruction.
    Preserves safe zone, rebuilds content sections using template DNA.
    """
    try:
        from docx import Document
        from docx.shared import RGBColor, Pt
        
        doc = Document(str(template_path))
        mappings = mapping.mappings
        template_dna = analysis.template_dna
        
        if not template_dna:
            logger.error("No template DNA found - cannot reconstruct")
            raise RenderError("Template DNA not extracted")
        
        logger.info(f"Rendering DOCX using template reconstruction")
        logger.info(f"  Template DNA: heading='{template_dna.heading_style_name}', body='{template_dna.body_style_name}'")
        logger.info(f"  Safe zone: paras 0-{template_dna.safe_zone_end_idx}")
        logger.info(f"  First content section: para {template_dna.first_content_section_idx}")
        
        # Parse the AI response - expecting array of sections
        sections_data = mappings.get("sections", [])
        if not sections_data:
            # Fallback: check if it's the old format
            if "sec_all" in mappings:
                logger.warning("Received old format, converting to sections")
                sections_data = [{"title": "Document", "body": mappings["sec_all"]}]
            else:
                logger.error("No sections data found in mappings")
                raise RenderError("No sections found in AI response")
        
        logger.info(f"  Rebuilding with {len(sections_data)} source sections")
        
        # Extract section titles for TOC update
        section_titles = [sec.get("title", f"Section {i+1}") for i, sec in enumerate(sections_data)]
        
        # Use first section title as document title for cover page
        document_title = section_titles[0] if section_titles else "Document"
        
        # Step 0: Update safe zone (cover page title + TOC) with new section titles
        try:
            _update_safe_zone(doc, template_dna, section_titles, document_title)
        except Exception as e:
            logger.warning(f"  Failed to update safe zone: {e}")
        
        # Step 1: Delete all paragraphs after safe zone
        paragraphs_to_remove = list(range(template_dna.safe_zone_end_idx, len(doc.paragraphs)))
        logger.info(f"  Removing {len(paragraphs_to_remove)} old content paragraphs")
        
        for idx in reversed(paragraphs_to_remove):
            if idx < len(doc.paragraphs):
                p = doc.paragraphs[idx]._element
                p.getparent().remove(p)
        
        # Step 2: Rebuild sections using template DNA
        for i, section in enumerate(sections_data):
            title = section.get("title", f"Section {i+1}")
            body = section.get("body", [])
            
            # Handle backwards compatibility - if body is a string, convert to array
            if isinstance(body, str):
                body = [{"type": "text", "content": body}]
            
            # DEDUPLICATION: Remove items with duplicate content
            # Strategy: If duplicates exist, prioritize "bullet" > "subheading" > "text"
            
            # 1. Group items by normalized content
            content_groups = {}
            for item in body:
                content = item.get("content", "").strip()
                if not content:
                    continue
                    
                # Normalize (remove bullets, numbers, case)
                normalized = re.sub(r'^[-*•\d.)\s]+', '', content).lower()
                
                if normalized not in content_groups:
                    content_groups[normalized] = []
                content_groups[normalized].append(item)
            
            # 2. Select best item for each unique content
            deduplicated_body = []
            
            # We want to preserve order, so we iterate through original body
            seen_normalized = set()
            
            for item in body:
                content = item.get("content", "").strip()
                normalized = re.sub(r'^[-*•\d.)\s]+', '', content).lower()
                
                if not normalized or normalized in seen_normalized:
                    continue
                
                # Get all versions of this content
                candidates = content_groups.get(normalized, [])
                
                # Pick the best one: Bullet > Subheading > Text
                best_item = item # Default to current
                
                has_bullet = any(c.get("type") == "bullet" for c in candidates)
                has_subheading = any(c.get("type") == "subheading" for c in candidates)
                
                if has_bullet:
                    # Find the first bullet version
                    best_item = next(c for c in candidates if c.get("type") == "bullet")
                elif has_subheading:
                    # Find the first subheading version
                    best_item = next(c for c in candidates if c.get("type") == "subheading")
                    
                deduplicated_body.append(best_item)
                seen_normalized.add(normalized)
            
            body = deduplicated_body
            logger.info(f"  Adding section {i+1}: '{title[:40]}' ({len(body)} body items)")
            
            # Add heading with template style
            heading_para = doc.add_paragraph(title, style=template_dna.heading_style_name)
            
            # Apply heading font properties
            if heading_para.runs:
                heading_run = heading_para.runs[0]
                if template_dna.heading_font_name:
                    heading_run.font.name = template_dna.heading_font_name
                if template_dna.heading_font_size:
                    from docx.shared import Pt
                    heading_run.font.size = Pt(template_dna.heading_font_size)
                if template_dna.heading_font_bold is not None:
                    heading_run.font.bold = template_dna.heading_font_bold
                if template_dna.heading_font_color:
                    try:
                        from docx.shared import RGBColor
                        heading_run.font.color.rgb = template_dna.heading_font_color
                    except:
                        pass
            
            # Add body items based on type
            for item in body:
                item_type = item.get("type", "text")
                content = item.get("content", "")
                
                if not content:
                    continue
                
                if item_type == "subheading":
                    # Try to use subheading style, fall back to body style with bold
                    try:
                        para = doc.add_paragraph(content, style=template_dna.subheading_style_name)
                    except:
                        logger.warning(f"Subheading style '{template_dna.subheading_style_name}' not found, using body style")
                        para = doc.add_paragraph(content, style=template_dna.body_style_name)
                        if para.runs:
                            para.runs[0].font.bold = True
                    
                    # Apply formatting
                    if para.runs and template_dna.subheading_font_name:
                        para.runs[0].font.name = template_dna.subheading_font_name
                    if para.runs and template_dna.subheading_font_size:
                        from docx.shared import Pt
                        para.runs[0].font.size = Pt(template_dna.subheading_font_size)
                    elif para.runs and template_dna.body_font_size:
                        # Fall back to slightly larger body size
                        from docx.shared import Pt
                        para.runs[0].font.size = Pt(int(template_dna.body_font_size * 1.1))
                    if para.runs and template_dna.subheading_font_bold is not None:
                        para.runs[0].font.bold = template_dna.subheading_font_bold
                        
                elif item_type == "bullet":
                    # Clean bullet text to remove any existing markers
                    bullet_content = clean_bullet_text(content)
                    
                    # Check if bullet style exists BEFORE trying to use it
                    has_bullet_style = False
                    try:
                        # Check if style exists without adding paragraph
                        _ = doc.styles[template_dna.bullet_style_name]
                        has_bullet_style = True
                    except:
                        has_bullet_style = False
                    
                    if has_bullet_style:
                        # Style exists, use it
                        para = doc.add_paragraph(bullet_content, style=template_dna.bullet_style_name)
                    else:
                        # Style doesn't exist, use manual formatting
                        logger.warning(f"Bullet style '{template_dna.bullet_style_name}' not found, using manual formatting")
                        para = doc.add_paragraph(style=template_dna.body_style_name)
                        # Add bullet manually
                        para.add_run("• " + bullet_content)
                        # Indent for bullet
                        from docx.shared import Pt
                        para.paragraph_format.left_indent = Pt(18)
                        para.paragraph_format.first_line_indent = Pt(-18)
                    
                    # Apply font formatting
                    if para.runs and template_dna.bullet_font_name:
                        para.runs[0].font.name = template_dna.bullet_font_name
                    elif para.runs and template_dna.body_font_name:
                        para.runs[0].font.name = template_dna.body_font_name
                    if para.runs and template_dna.bullet_font_size:
                        from docx.shared import Pt
                        para.runs[0].font.size = Pt(template_dna.bullet_font_size)
                    elif para.runs and template_dna.body_font_size:
                        from docx.shared import Pt
                        para.runs[0].font.size = Pt(template_dna.body_font_size)
                        
                else:  # text
                    para = doc.add_paragraph(content, style=template_dna.body_style_name)
                    if para.runs:
                        if template_dna.body_font_name:
                            para.runs[0].font.name = template_dna.body_font_name
                        if template_dna.body_font_size:
                            from docx.shared import Pt
                            para.runs[0].font.size = Pt(template_dna.body_font_size)
                        if template_dna.body_font_bold is not None:
                            para.runs[0].font.bold = template_dna.body_font_bold
                        if template_dna.body_font_italic is not None:
                            para.runs[0].font.italic = template_dna.body_font_italic
            
            # Add some spacing between sections
            if i < len(sections_data) - 1:
                doc.add_paragraph()  # Empty paragraph for spacing
        
        # Save
        output_path.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(output_path))
        logger.info(f"Saved reconstructed document to: {output_path}")
        return output_path
        
    except Exception as e:
        logger.error(f"Rendering failed: {e}")
        raise RenderingError(f"Failed to render DOCX: {str(e)}", details=str(e))


def _render_pptx_sections(
    template_path: Path,
    output_path: Path,
    mapping: SectionMapping,
    analysis: TemplateAnalysis
) -> Path:
    """
    Render PPTX template with section-based content replacement.
    Each slide = one section. Replace body text while keeping title and layout.
    """
    try:
        from pptx import Presentation
        
        prs = Presentation(str(template_path))
        mappings = mapping.mappings
        
        for section in analysis.sections:
            section_id = section.section_id
            new_content = mappings.get(section_id, "")
            
            if not new_content:
                continue
            
            # Extract slide index from section_id (e.g., "slide_0" -> 0)
            if not section_id.startswith("slide_"):
                continue
            
            try:
                slide_idx = int(section_id.split("_")[1])
            except (ValueError, IndexError):
                continue
            
            if slide_idx >= len(prs.slides):
                continue
            
            slide = prs.slides[slide_idx]
            
            # Find body shapes (non-title text frames)
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                
                # Skip title shapes
                is_title = False
                if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                    ph_type = shape.placeholder_format.type
                    if ph_type in [1, 2, 3]:
                        is_title = True
                
                if not is_title:
                    # Replace body text
                    tf = shape.text_frame
                    if tf.paragraphs:
                        # Get style from first paragraph
                        first_para = tf.paragraphs[0]
                        
                        # Set content in first paragraph
                        if first_para.runs:
                            first_para.runs[0].text = new_content
                            for run in first_para.runs[1:]:
                                run.text = ""
                        else:
                            first_para.text = new_content
                        
                        # Clear other paragraphs
                        for para in tf.paragraphs[1:]:
                            for run in para.runs:
                                run.text = ""
                    
                    # Only replace first body shape per slide
                    break
        
        # Ensure output directory exists
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        prs.save(str(output_path))
        return output_path
        
    except Exception as e:
        raise RenderingError(f"Failed to render PPTX: {str(e)}", details=str(e))
