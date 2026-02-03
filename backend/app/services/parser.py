"""
Content extraction from source documents (DOCX, PDF, PPTX).
Extracts text blocks preserving order but ignoring original formatting.
"""
import re
from pathlib import Path
from typing import Literal

from pydantic import BaseModel

from app.core.exceptions import ParsingError, UnsupportedFileTypeError


class ContentBlock(BaseModel):
    """A block of extracted content."""
    id: str
    type: Literal["heading", "paragraph", "list", "table"]
    content: str


class ExtractedContent(BaseModel):
    """Container for all extracted content blocks."""
    blocks: list[ContentBlock]
    source_file: str


def extract_content(file_path: Path) -> ExtractedContent:
    """
    Extract content blocks from a document.
    
    Args:
        file_path: Path to the source document
        
    Returns:
        ExtractedContent with ordered content blocks
        
    Raises:
        UnsupportedFileTypeError: If file type is not supported
        ParsingError: If parsing fails
    """
    suffix = file_path.suffix.lower()
    
    if suffix == ".docx":
        return _extract_from_docx(file_path)
    elif suffix == ".pdf":
        return _extract_from_pdf(file_path)
    elif suffix == ".pptx":
        return _extract_from_pptx(file_path)
    else:
        raise UnsupportedFileTypeError(f"Unsupported file type: {suffix}")


def _extract_from_docx(file_path: Path) -> ExtractedContent:
    """Extract content from DOCX file."""
    try:
        from docx import Document
        from docx.oxml.ns import qn
        
        doc = Document(str(file_path))
        blocks: list[ContentBlock] = []
        block_id = 0
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
                
            # Determine block type based on style
            style_name = para.style.name.lower() if para.style else ""
            
            if "heading" in style_name or "title" in style_name:
                block_type = "heading"
            elif "list" in style_name or para.text.strip().startswith(("-", "•", "*")):
                block_type = "list"
            else:
                block_type = "paragraph"
            
            blocks.append(ContentBlock(
                id=f"b{block_id}",
                type=block_type,
                content=text
            ))
            block_id += 1
        
        # Extract table content
        for table in doc.tables:
            table_text_parts = []
            for row in table.rows:
                row_cells = [cell.text.strip() for cell in row.cells]
                if any(row_cells):
                    table_text_parts.append(" | ".join(row_cells))
            
            if table_text_parts:
                blocks.append(ContentBlock(
                    id=f"b{block_id}",
                    type="table",
                    content="\n".join(table_text_parts)
                ))
                block_id += 1
        
        return ExtractedContent(blocks=blocks, source_file=file_path.name)
        
    except Exception as e:
        raise ParsingError(f"Failed to parse DOCX: {str(e)}", details=str(e))


def _extract_from_pdf(file_path: Path) -> ExtractedContent:
    """Extract content from PDF file."""
    try:
        import fitz  # PyMuPDF
        
        doc = fitz.open(str(file_path))
        blocks: list[ContentBlock] = []
        block_id = 0
        
        for page in doc:
            # Extract text blocks with their positions
            text_blocks = page.get_text("blocks")
            
            for block in text_blocks:
                # block format: (x0, y0, x1, y1, text, block_no, block_type)
                if len(block) >= 5:
                    text = block[4].strip()
                    if text and not text.startswith("<image"):
                        # Simple heuristic: short lines at top might be headings
                        block_type = "paragraph"
                        if len(text) < 100 and text.isupper():
                            block_type = "heading"
                        elif text.startswith(("-", "•", "*", "►")):
                            block_type = "list"
                        
                        blocks.append(ContentBlock(
                            id=f"b{block_id}",
                            type=block_type,
                            content=text
                        ))
                        block_id += 1
        
        doc.close()
        return ExtractedContent(blocks=blocks, source_file=file_path.name)
        
    except Exception as e:
        raise ParsingError(f"Failed to parse PDF: {str(e)}", details=str(e))


def _extract_from_pptx(file_path: Path) -> ExtractedContent:
    """Extract content from PPTX file."""
    try:
        from pptx import Presentation
        
        prs = Presentation(str(file_path))
        blocks: list[ContentBlock] = []
        block_id = 0
        
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                    
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if not text:
                        continue
                    
                    # Detect block type
                    # Title shapes usually have placeholder type
                    block_type = "paragraph"
                    if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                        ph_type = shape.placeholder_format.type
                        # Type 1 = Title, Type 2 = Center Title
                        if ph_type in [1, 2, 3]:
                            block_type = "heading"
                    elif paragraph.level > 0:
                        block_type = "list"
                    
                    blocks.append(ContentBlock(
                        id=f"b{block_id}",
                        type=block_type,
                        content=text
                    ))
                    block_id += 1
            
            # Extract table content from slides
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    table_text_parts = []
                    for row in table.rows:
                        row_cells = [cell.text.strip() for cell in row.cells]
                        if any(row_cells):
                            table_text_parts.append(" | ".join(row_cells))
                    
                    if table_text_parts:
                        blocks.append(ContentBlock(
                            id=f"b{block_id}",
                            type="table",
                            content="\n".join(table_text_parts)
                        ))
                        block_id += 1
        
        return ExtractedContent(blocks=blocks, source_file=file_path.name)
        
    except Exception as e:
        raise ParsingError(f"Failed to parse PPTX: {str(e)}", details=str(e))


def content_to_text_summary(content: ExtractedContent) -> str:
    """
    Convert extracted content to a text summary for AI processing.
    
    Args:
        content: Extracted content blocks
        
    Returns:
        Formatted text summary
    """
    lines = []
    for block in content.blocks:
        prefix = f"[{block.type.upper()}]"
        lines.append(f"{prefix} {block.content}")
    
    return "\n\n".join(lines)
